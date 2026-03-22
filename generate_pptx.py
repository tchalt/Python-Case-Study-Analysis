from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import os

def set_cell_text(cell, text, bold=False, font_size=10, font_color=None, alignment=PP_ALIGN.CENTER):
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if font_color:
            paragraph.font.color.rgb = RGBColor(*font_color)
        paragraph.alignment = alignment

def apply_table_style(table):
    """Apply blue header and alternating row colors to match template images"""
    # Header row
    for c in range(len(table.columns)):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(73, 124, 182) # Template blue
        set_cell_text(cell, cell.text, bold=True, font_color=(255, 255, 255))

    # Data rows
    for r in range(1, len(table.rows)):
        fill_color = RGBColor(217, 226, 243) if r % 2 == 0 else RGBColor(255, 255, 255)
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            set_cell_text(cell, cell.text, font_size=9)

def create_presentation():
    # Load data
    df = pd.read_csv('financial_data.csv')
    y5 = df.iloc[-1]
    
    prs = Presentation()

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Lithuania PV Production Financial Plan"
    subtitle.text = "Market-Scaled Strategic Investment Evaluation (250 MW)\nMarch 2026"

    # Slide 2: Baseline vs With-Project Summary
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "1. Baseline vs With-Project summary"
    rows, cols = 5, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    headers = ['Metric', 'Baseline (Outsource)', 'With-Project (Local)']
    for c, h in enumerate(headers): table.cell(0, c).text = h
    
    data = [
        ['Supply Chain Risk', 'High (Long lead times)', 'Low (Local control)'],
        ['Production Cost', 'Market Price + Margin', 'Internal Variable Cost'],
        ['Strategic Value', 'Dependency on Imports', 'EU "Made in EU" Premium'],
        ['CIT Rate', '15%', '0% (Klaipeda FEZ)']
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row): table.cell(r+1, c).text = val
    apply_table_style(table)

    # Slide 3: Capacity Model (Template 1)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "TABLE TEMPLATE 1 — Capacity & ramp-up plan"
    rows, cols = 6, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    headers = ['Year', 'Max capacity', 'Utilization %', 'Scrap %', 'Effective capacity']
    for c, h in enumerate(headers): table.cell(0, c).text = h
    
    for r, row in df.iterrows():
        data_row = [
            f"Year {int(r+1)}",
            f"{row['Max_Capacity_MW']:.0f}",
            f"{row['Utilization_%']*100:.1f}%",
            f"{row['Scrap_%']*100:.1f}%",
            f"{row['Effective_Capacity_MW']:.2f}"
        ]
        for c, val in enumerate(data_row): table.cell(r+1, c).text = val
    apply_table_style(table)

    # Slide 4: Sales Model (Template 2)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "TABLE TEMPLATE 2 — Sales forecast"
    rows, cols = 6, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    headers = ['Year', 'Units sold', 'Price (EUR/unit)', 'Revenue (EUR)', 'Justification (short)']
    for c, h in enumerate(headers): table.cell(0, c).text = h
    
    justifications = ['Market entry', 'Growth phase', 'Scale reach', 'Maturity', 'Full capacity']
    for r, row in df.iterrows():
        data_row = [
            f"Year {int(r+1)}",
            f"{row['Effective_Capacity_MW']:.2f}",
            f"{row['Price_EUR_MW']:,.0f}",
            f"{row['Revenue_EUR']:,.0f}",
            justifications[r]
        ]
        for c, val in enumerate(data_row): table.cell(r+1, c).text = val
    apply_table_style(table)

    # Slide 5: Variable Cost Breakdown (Template 3)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "TABLE TEMPLATE 3 — Variable cost per unit breakdown"
    rows, cols = 6, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    headers = ['Component', 'Driver', 'Formula', 'EUR/unit', 'Source']
    for c, h in enumerate(headers): table.cell(0, c).text = h
    
    vcost_data = [
        ['Materials/ingredients', 'LONGi HPBC 2.0', 'Direct', '80,700', 'LONGi 2025 Spec'],
        ['Packaging', 'Export grade', 'Fixed', '1,500', 'Logistics Spec'],
        ['Direct labor', 'Lithuania Min', '800h * 16.5', '13,200', 'VDI 2026 Std'],
        ['Energy (if relevant)', 'Industrial rate', '0.18/kWh', '2,500', 'Eurostat'],
        ['Platform fee / other', 'Scrap Effect', 'Cost/(1-S)-C', f"{y5['Scrap_Effect_EUR_MW']:,.0f}", 'Industry Std']
    ]
    for r, row in enumerate(vcost_data):
        for c, val in enumerate(row): table.cell(r+1, c).text = val
    apply_table_style(table)

    # Slide 6: Incremental Fixed Costs (Template 4)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "TABLE TEMPLATE 4 — Incremental fixed costs"
    rows, cols = 6, 6
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    headers = ['Year', 'Maintenance', 'Extra staff', 'Rent/overhead', 'Other', 'Total fixed']
    for c, h in enumerate(headers): table.cell(0, c).text = h
    
    for r in range(5):
        data_row = [f"Year {r+1}", '432,000', '288,000', '96,000', '384,000', '1,200,000']
        for c, val in enumerate(data_row): table.cell(r+1, c).text = val
    apply_table_style(table)

    # Slide 7: Operating Statement (Template 5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "TABLE TEMPLATE 5 — Project operating statement"
    rows, cols = 6, 7
    table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1.5), Inches(9.6), Inches(3)).table
    headers = ['Year', 'Units', 'Revenue', 'Variable costs', 'Contribution', 'Fixed costs', 'EBITDA']
    for c, h in enumerate(headers): table.cell(0, c).text = h
    
    for r, row in df.iterrows():
        data_row = [
            f"Year {int(r+1)}",
            f"{row['Effective_Capacity_MW']:.1f}",
            f"{row['Revenue_EUR']:,.0f}",
            f"{row['Total_Var_Costs_EUR']:,.0f}",
            f"{row['Contribution_Margin_EUR']:,.0f}",
            '1,200,000',
            f"{row['EBITDA_EUR']:,.0f}"
        ]
        for c, val in enumerate(data_row): table.cell(r+1, c).text = val
    apply_table_style(table)

    # Slide 8: Break-even Analysis (Template 6)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "TABLE TEMPLATE 6 — Break-even"
    q_be = y5['Fixed_Costs_EUR'] / (y5['Price_EUR_MW'] - y5['Total_Var_Cost_Per_MW'])
    mos = (y5['Effective_Capacity_MW'] - q_be) / y5['Effective_Capacity_MW'] * 100
    
    rows, cols = 8, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table
    headers = ['Item', 'Value', 'Notes']
    for c, h in enumerate(headers): table.cell(0, c).text = h
    
    be_data = [
        ['Price (EUR/unit)', f"{y5['Price_EUR_MW']:,.0f}", 'EUR/MW'],
        ['Variable cost (EUR/unit)', f"{y5['Total_Var_Cost_Per_MW']:,.0f}", 'EUR/MW'],
        ['Contribution (EUR/unit)', f"{(y5['Price_EUR_MW'] - y5['Total_Var_Cost_Per_MW']):,.0f}", 'Price - Var cost'],
        ['Fixed costs (EUR/period)', '1,200,000', 'Annual burden'],
        ['Break-even units (Q_BE)', f"{q_be:,.2f}", 'Fixed / Contrib unit'],
        ['Break-even revenue', f"{q_be * y5['Price_EUR_MW']:,.0f}", 'Q_BE x Price'],
        ['Margin of safety', f"{mos:.1f}%", '(Exp - Q_BE)/Exp']
    ]
    for r, row in enumerate(be_data):
        for c, val in enumerate(row): table.cell(r+1, c).text = val
    apply_table_style(table)

    # Slide 9: Assumptions & Sources (Template 7)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "TABLE TEMPLATE 7 — Assumptions & Sources"
    rows, cols = 8, 6
    table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1.5), Inches(9.6), Inches(4)).table
    headers = ['Assumption', 'Value', 'Unit', 'Source', 'Why reasonable', 'Low/Base/High']
    for c, h in enumerate(headers): table.cell(0, c).text = h
    
    ass_data = [
        ['Price', '126,000', 'EUR/MW', 'LONGi 2025', 'HPBC 2.0 premium', 'Base'],
        ['Volume / demand', '250', 'MW', 'Market study', 'Realistic domestic share', 'Base'],
        ['Ramp-up (utilization)', '50%-100%', '%', 'Project plan', 'Staged growth', 'Base'],
        ['Wage rate', '16.50', 'EUR/h', 'VDI 2026', 'Lithuanian standard', 'Base'],
        ['Energy price', '0.18', 'EUR/kWh', 'Eurostat', 'Baltic industrial', 'Base'],
        ['Scrap / yield', '2% - 5%', '%', 'Industry Std', 'Improving ramp-up', 'Base'],
        ['Maintenance', '432,000', 'EUR/yr', 'Vendor spec', 'Scaled for 250MW', 'Base']
    ]
    for r, row in enumerate(ass_data):
        for c, val in enumerate(row): table.cell(r+1, c).text = val
    apply_table_style(table)

    # Slide 10: Charts
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Financial Performance Visualizations"
    if os.path.exists('volume_rampup.png'):
        slide.shapes.add_picture('volume_rampup.png', Inches(0.5), Inches(1.5), height=Inches(3.5))
    if os.path.exists('breakeven_chart.png'):
        slide.shapes.add_picture('breakeven_chart.png', Inches(5), Inches(1.5), height=Inches(3.5))

    prs.save('Production_Financial_Plan_Presentation.pptx')
    print("Reordered and updated Presentation generated: Production_Financial_Plan_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
